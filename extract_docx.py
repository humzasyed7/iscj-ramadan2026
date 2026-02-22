#!/usr/bin/env python3

"""Aggregate source material and the active site brief into a single bundle.

The generated markdown file can be shared with the assistant to provide
comprehensive context for building a website based on user-provided assets and
requirements.
"""

from __future__ import annotations

import argparse
from datetime import UTC, datetime
from pathlib import Path
from typing import Callable, Dict, List, Tuple


def read_docx(file_path: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except ModuleNotFoundError as exc:  # pragma: no cover
        raise RuntimeError(
            "python-docx is required to read .docx files. Install it via 'pip install python-docx'."
        ) from exc

    document = Document(file_path)
    pieces: List[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            pieces.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                pieces.append(" | ".join(cells))

    return "\n".join(pieces)


def read_pdf(file_path: Path) -> str:
    reader_cls = None
    try:  # pragma: no cover - import guarding is not easily testable here
        from pypdf import PdfReader  # type: ignore

        reader_cls = PdfReader
    except ModuleNotFoundError:
        try:
            from PyPDF2 import PdfReader  # type: ignore

            reader_cls = PdfReader
        except ModuleNotFoundError as exc:
            raise RuntimeError(
                "pypdf or PyPDF2 is required to read .pdf files. Install via 'pip install pypdf'."
            ) from exc

    reader = reader_cls(str(file_path))  # type: ignore[call-arg]
    pages: List[str] = []
    for page in reader.pages:  # type: ignore[attr-defined]
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(text)

    return "\n\n".join(pages)


def read_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ModuleNotFoundError as exc:  # pragma: no cover
        raise RuntimeError(
            "python-pptx is required to read .pptx files. Install it via 'pip install python-pptx'."
        ) from exc

    presentation = Presentation(file_path)
    slides_output: List[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        fragments: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    fragments.append(text)
        if fragments:
            slides_output.append(f"Slide {index}:\n" + "\n".join(fragments))

    return "\n\n".join(slides_output)


def read_plain_text(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore").strip()


READERS: Dict[str, Callable[[Path], str]] = {
    ".docx": read_docx,
    ".pdf": read_pdf,
    ".pptx": read_pptx,
    ".txt": read_plain_text,
    ".md": read_plain_text,
}


def sanitize_heading(text: str) -> str:
    return text.replace("`", "\u2019")


def load_site_brief(path: Path) -> str:
    if not path.exists():
        return ""
    return path.read_text(encoding="utf-8").strip()


def gather_sources(source_dir: Path) -> Tuple[List[Tuple[Path, str]], List[Path], List[Tuple[Path, Exception]]]:
    gathered: List[Tuple[Path, str]] = []
    unsupported: List[Path] = []
    failures: List[Tuple[Path, Exception]] = []

    for file_path in sorted(source_dir.glob("*")):
        if file_path.is_dir():
            continue

        reader = READERS.get(file_path.suffix.lower())
        if not reader:
            unsupported.append(file_path)
            continue

        try:
            text = reader(file_path).strip()
            gathered.append((file_path, text))
        except Exception as exc:  # pragma: no cover - defensive logging
            failures.append((file_path, exc))

    return gathered, unsupported, failures


def build_bundle(request_text: str, sources: List[Tuple[Path, str]]) -> str:
    lines: List[str] = []
    timestamp = datetime.now(UTC).isoformat(timespec="seconds")
    lines.append("# Website Request Bundle")
    lines.append(f"Generated: {timestamp}")
    lines.append("")

    if request_text:
        lines.append("## Site Brief (`request/site-brief.yaml`)")
        lines.append("```yaml")
        lines.append(request_text)
        lines.append("```")
        lines.append("")
    else:
        lines.append("## Site Brief")
        lines.append("_No site brief found. Fill out `request/site-brief.yaml` to capture your requirements._")
        lines.append("")

    lines.append("## Source Material")
    if not sources:
        lines.append("_No supported source files discovered in the `source/` directory._")
    else:
        for file_path, text in sources:
            heading = sanitize_heading(file_path.name)
            lines.append(f"### {heading}")
            if text:
                lines.append("```text")
                lines.append(text)
                lines.append("```")
            else:
                lines.append("_No extractable text content found._")
            lines.append("")

    return "\n".join(lines).strip() + "\n"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Aggregate the active site brief and source documents into a single markdown file."
    )
    parser.add_argument(
        "--source-dir",
        default="source",
        type=Path,
        help="Directory containing source material to ingest.",
    )
    parser.add_argument(
        "--request-file",
        default=Path("request/site-brief.yaml"),
        type=Path,
        help="YAML or text file describing the desired website.",
    )
    parser.add_argument(
        "--output",
        default=Path("build/context-bundle.md"),
        type=Path,
        help="Destination markdown file that consolidates the inputs.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    source_dir: Path = args.source_dir
    if not source_dir.exists():  # pragma: no cover
        raise SystemExit(f"Source directory not found: {source_dir}")

    request_text = load_site_brief(Path(args.request_file))
    gathered, unsupported, failures = gather_sources(source_dir)

    bundle = build_bundle(request_text, gathered)

    output_path: Path = args.output
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(bundle, encoding="utf-8")

    print(f"✅ Context bundle written to {output_path.resolve()}")

    if unsupported:
        print("⚠️  Unsupported files (skipped):")
        for path in unsupported:
            print(f"   • {path.name}")

    if failures:
        print("❌ Files that could not be processed:")
        for path, error in failures:
            print(f"   • {path.name}: {error}")


if __name__ == "__main__":
    main()
